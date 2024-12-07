"""
goldenbell.py    골든볼 파워포인트  Ver 1.0_240518
"""
from collections import abc
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN

# 프레젠테이션 만들기
prs = Presentation()

full_width = prs.slide_width
full_height = prs.slide_height
font_size = Pt(48)

quiz = {
    '구석기 시대에 돌을 깨뜨리거나 떼내어 만든 도구는?': '뗀석기',
    '삼국 중 가장 먼저 한강 유역을 차지한 나라는?': '백제',
    '왕건이 후삼국을 통일하고 세운 나라는?': '고려',
    '조선 세종시대 노비 출신 과학자로, 자격루를 만든 사람은?': '장영실',
    '조선 후기 백성이 잘 살고 나라가 강해지는 법을 연구한 학문은?': '실학',
}

for key, value in quiz.items():
    for i in range(2):
        # 빈 슬라이드 추가
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 슬라이드에 텍스트박스 추가
        textbox = slide.shapes.add_textbox(
            left=0,
            top=full_height / 2 - font_size,
            width=full_width,
            height=0
        )
        textbox.text_frame.word_wrap = True

        # 텍스트박스의 문단 설정
        paragraph = textbox.text_frame.paragraphs[0]
        paragraph.font.bold = True
        paragraph.font.size = font_size
        paragraph.alignment = PP_ALIGN.CENTER

        # 문단에 텍스트 추가
        if i == 0:
            paragraph.text = key
        else:
            paragraph.text = value

prs.save("./school/역사골든벨.pptx")
