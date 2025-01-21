"""
name_card.py 행사 명찰 제작  Ver 1.0_240508
"""
# 01. 작업 파일 불러 오기
import copy
import pandas as pd            # pandas 모듈 불러오기
from pptx import Presentation  # Presentation 모듈 불러오기
from pptx.util import Pt       # 폰트 크기 조정을 위해 로드

df = pd.read_excel("./data/행사 참석자 명단.xlsx")   # 참석자 명단 불러와 df에 저장
prs = Presentation("./data/명찰 양식.pptx")          # 명찰 양식 불러와 prs 객체 생성

# 02. 사용자 함수 정의
def duplicate_slide(prs, org_slide):
    # add_slide() 메소드를 이용해 새로운 슬라이드 생성한 후 개체를 복사해 옴
    copied_slide = prs.slides.add_slide(org_slide.slide_layout)
    for shape in org_slide.shapes:
        org_elment = shape.element
        new_element = copy.deepcopy(org_elment)
        copied_slide.shapes._spTree.insert_element_before(new_element, "p:extLst")

    for value in org_slide.part.rels.values():
        if "notesSlide" not in value.reltype:
            copied_slide.part.rels.get_or_add(
                value.reltype,
                value._target
            )
    return copied_slide

# 03. 명찰 제작하기
nametag_count = 0
for person_count in range(len(df)):  # 참석자 수(32명)만큼 반복문 실행
    # 슬라이드 1장 당 4개의 명찰이 있으므로
    # 현재 반복문이 실행되고 있는 person_count를 4로 나누고
    # 남은 값이 0일 때(=4의 배수일 때) 슬라이드를 복사하고 데이터를 입력함
    if person_count % 4 == 0:
        new_slide = duplicate_slide(prs, prs.slides[0])  # 첫 번째 슬라이드를 복사하여 추가
        for shape in new_slide.shapes:                   # 슬라이드에 있는 개체를 체크하여 데이터 입력
            # 명찰 입력이 완료된 개수(nametag_count)가
            # 명단 데이터 개수(len(df))보다 적을 때만 실행
            if nametag_count < len(df):
                if shape.shape_type == 17 and shape.text == "소속":
                    shape.text_frame.paragraphs[0].text = df.iloc[nametag_count, 1]
                    shape.text_frame.paragraphs[0].font.size = Pt(24)
                    shape.text_frame.paragraphs[0].font.bold = True
                    shape.text_frame.paragraphs[0].font.name = "맑은 고딕"
                elif shape.shape_type == 17 and shape.text == "이름":
                    shape.text_frame.paragraphs[0].text = df.iloc[nametag_count, 0]
                    shape.text_frame.paragraphs[0].font.size = Pt(60)
                    shape.text_frame.paragraphs[0].font.bold = True
                    shape.text_frame.paragraphs[0].font.name = "맑은 고딕"
                    # 이름 입력까지 완료되면 입력 완료된 개수(nametag_count) 업데이트
                    nametag_count += 1

prs.save("./output/명찰 자동제작 결과.pptx")  # 새로운 파일로 저장